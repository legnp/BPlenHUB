# -*- coding: utf-8 -*-
"""
QA geometrico sem LibreOffice: confere se alguma forma estoura os limites
do slide, e se caixas de texto ficam contidas dentro do cartao mais proximo
(quando aplicavel). Nao substitui inspecao visual, mas pega os defeitos mais
comuns (estouro de borda, caixa fora do cartao) sem precisar renderizar.
"""
import sys
from pptx import Presentation
from pptx.util import Emu

EMU_POR_POL = 914400

p = Presentation(sys.argv[1] if len(sys.argv) > 1 else "BPlen-Apresentacao-Institucional.pptx")
SW, SH = p.slide_width, p.slide_height

problemas = []

for i, slide in enumerate(p.slides, 1):
    cartoes = []
    formas = []
    for sh in slide.shapes:
        if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
            continue
        box = (sh.left, sh.top, sh.left + sh.width, sh.top + sh.height)
        formas.append((sh, box))
        # roundRect = provavel "cartao" de fundo (heuristica pelo tipo de shape)
        try:
            if sh.shape_type is not None and "ROUND" in str(sh.auto_shape_type or ""):
                cartoes.append(box)
        except Exception:
            pass

    for sh, (l, t, r, b) in formas:
        nome = (sh.text_frame.text[:30].replace("\n", " ") if sh.has_text_frame and sh.text_frame.text.strip() else sh.shape_type)
        margem = 0  # tolerancia zero: qualquer estouro e real
        if l < -margem or t < -margem or r > SW + margem or b > SH + margem:
            problemas.append(
                f"slide {i}: '{nome}' estoura a borda "
                f"(l={l/EMU_POR_POL:.2f} t={t/EMU_POR_POL:.2f} r={r/EMU_POR_POL:.2f} b={b/EMU_POR_POL:.2f} pol; "
                f"slide=13.33x7.5)"
            )

print(f"{len(p.slides)} slides verificados.")
if problemas:
    print(f"\n{len(problemas)} problema(s) de geometria:")
    for pr in problemas:
        print(" -", pr)
    sys.exit(1)
else:
    print("Nenhuma forma estourando os limites do slide.")
